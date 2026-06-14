from __future__ import annotations

import io
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def export_outline_to_pptx(outline: dict[str, Any]) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_text = str(outline.get("title") or "")
    chapters = outline.get("chapters") if isinstance(outline.get("chapters"), list) else []
    slides = outline.get("slides") if isinstance(outline.get("slides"), list) else []

    _add_title_slide(prs, title_text)

    for ch in chapters:
        if not isinstance(ch, dict):
            continue
        ch_id = str(ch.get("chapter_id") or "")
        ch_title = str(ch.get("title") or "")
        _add_chapter_divider(prs, ch_title, ch_id)

        ch_slide_ids = set(ch.get("slide_ids") or [])
        for s in slides:
            if not isinstance(s, dict):
                continue
            if str(s.get("slide_id") or "") not in ch_slide_ids:
                continue
            _add_content_slide(prs, str(s.get("title") or ""), s)

    for s in slides:
        if not isinstance(s, dict):
            continue
        sid = str(s.get("slide_id") or "")
        if chapters and any(sid in set(ch.get("slide_ids") or []) for ch in chapters if isinstance(ch, dict)):
            continue
        _add_content_slide(prs, str(s.get("title") or ""), s)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _add_title_slide(prs: Presentation, title: str) -> None:
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    left = Inches(1.5)
    top = Inches(2.5)
    width = Inches(10)
    height = Inches(2)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x28, 0x64, 0xD8)
    p.alignment = PP_ALIGN.CENTER


def _add_chapter_divider(prs: Presentation, chapter_title: str, chapter_id: str) -> None:
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    left = Inches(1.5)
    top = Inches(3.0)
    width = Inches(10)
    height = Inches(1.5)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = chapter_title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x28, 0x64, 0xD8)
    p.alignment = PP_ALIGN.CENTER


def _add_content_slide(prs: Presentation, slide_title: str, slide_data: dict[str, Any]) -> None:
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    left = Inches(1.0)
    top = Inches(0.6)
    width = Inches(11.333)
    height = Inches(0.9)

    shape = slide.shapes.add_shape(
        1, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x28, 0x64, 0xD8)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = slide_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.LEFT

    # Key message
    key_msg = str(slide_data.get("key_message") or "")
    if key_msg:
        left = Inches(1.2)
        top = Inches(1.8)
        width = Inches(10.933)
        height = Inches(0.7)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf2 = txBox.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = key_msg
        p2.font.size = Pt(14)
        p2.font.italic = True
        p2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # Bullets
    bullets = slide_data.get("bullets") if isinstance(slide_data.get("bullets"), list) else []
    left = Inches(1.2)
    top = Inches(2.8)
    width = Inches(10.933)
    height = Inches(4.0)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf3 = txBox.text_frame
    tf3.word_wrap = True

    for idx, bullet in enumerate(bullets):
        if not isinstance(bullet, dict):
            continue
        text = str(bullet.get("text") or "")
        if not text.strip():
            continue

        if idx == 0:
            p3 = tf3.paragraphs[0]
        else:
            p3 = tf3.add_paragraph()

        p3.text = f"• {text}"
        p3.font.size = Pt(16)
        p3.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p3.space_after = Pt(10)
        p3.level = 0
